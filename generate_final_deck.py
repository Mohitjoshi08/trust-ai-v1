import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = r'C:\Users\mohit\OneDrive\Desktop\AIC\ps\AIC_Talent-Brand_PPT-Template (1).pptx'
OUTPUT = r'C:\Users\mohit\OneDrive\Desktop\AIC\Team_Trace_Trace.ai.pptx'

prs = Presentation(TEMPLATE)

# Add exactly 2 blank slides for the content
blank_layout = prs.slide_layouts[6]
s1 = prs.slides.add_slide(blank_layout) # Problem
s2 = prs.slides.add_slide(blank_layout) # Solution

# Delete all slides except the Cover (index 0) and the two we just added.
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
slides_to_keep = [slides[0], slides[-2], slides[-1]]
for sld in slides:
    if sld not in slides_to_keep:
        xml_slides.remove(sld)

W = 13.33; H = 7.50
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PUR = RGBColor(0x7B, 0x2F, 0xBE)
LGRAY = RGBColor(0xF5, 0xF5, 0xF5)
MGRAY = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY = RGBColor(0x50, 0x50, 0x50)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x7B, 0x06)

def I(v): return Inches(v)

def R(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(l), I(t), I(w), I(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, fs=10, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or DARK
    bx = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = 'Arial'
    run.font.color.rgb = color
    return bx

def Arrow(slide, l, t, w, h, fill=PUR):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, I(l), I(t), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def RightArrow(slide, l, t, w, h, fill=PUR):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, I(l), I(t), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

# ==========================================
# SLIDE 1: PROBLEM
# ==========================================
R(s1, 0, 0, W, 0.7, DARK)
T(s1, 0.2, 0.15, 12, 0.4, "A dashboard tells you WHAT changed. It rarely tells you WHY.", fs=22, bold=True, color=WHITE)

# KPI Card
R(s1, 0.5, 1.4, 3.5, 1.8, WHITE, border=MGRAY)
T(s1, 0.5, 1.7, 3.5, 0.4, "REVENUE", fs=20, bold=True, align=PP_ALIGN.CENTER, color=DGRAY)
T(s1, 0.5, 2.2, 3.5, 0.6, "↓ 32%", fs=44, bold=True, align=PP_ALIGN.CENTER, color=RED)

# Vertical flow in middle
cx = 4.8
bw = 3.8
cx_mid = cx + bw/2

R(s1, cx, 1.0, bw, 0.6, LGRAY, border=MGRAY)
T(s1, cx, 1.15, bw, 0.4, "Dashboard", fs=14, bold=True, align=PP_ALIGN.CENTER)

Arrow(s1, cx_mid - 0.1, 1.7, 0.2, 0.3, MGRAY)

R(s1, cx, 2.1, bw, 0.6, LGRAY, border=MGRAY)
T(s1, cx, 2.25, bw, 0.4, "“What caused it?”", fs=14, bold=True, align=PP_ALIGN.CENTER, color=PUR)

Arrow(s1, cx_mid - 0.1, 2.8, 0.2, 0.3, MGRAY)

R(s1, cx, 3.2, bw, 0.6, LGRAY, border=MGRAY)
T(s1, cx, 3.35, bw, 0.4, "SQL / BI dashboards", fs=14, bold=True, align=PP_ALIGN.CENTER)

Arrow(s1, cx_mid - 0.1, 3.9, 0.2, 0.3, MGRAY)

# Wide box
R(s1, cx - 0.5, 4.3, bw + 1.0, 0.6, LGRAY, border=MGRAY)
T(s1, cx - 0.5, 4.45, bw + 1.0, 0.4, "GitHub / Jira / Slack / support tickets", fs=13, bold=True, align=PP_ALIGN.CENTER)

Arrow(s1, cx_mid - 0.1, 5.0, 0.2, 0.3, MGRAY)

R(s1, cx, 5.4, bw, 0.6, LGRAY, border=MGRAY)
T(s1, cx, 5.55, bw, 0.4, "Manual investigation", fs=14, bold=True, align=PP_ALIGN.CENTER)

Arrow(s1, cx_mid - 0.1, 6.1, 0.2, 0.3, MGRAY)

R(s1, cx, 6.5, bw, 0.6, DARK)
T(s1, cx, 6.65, bw, 0.4, "Business decision", fs=14, bold=True, align=PP_ALIGN.CENTER, color=WHITE)

# Highlighted statement
R(s1, 9.2, 2.6, 3.7, 1.4, RGBColor(0xED, 0xE7, 0xF6), border=PUR)
T(s1, 9.4, 2.9, 3.3, 0.8, "“The real bottleneck isn't detecting a KPI change. It is translating the change into a defensible business decision.”", fs=14, bold=True, color=PUR, align=PP_ALIGN.CENTER)

# Contrast
T(s1, 0.5, 6.3, 4.0, 0.3, "TODAY", fs=14, bold=True, color=DGRAY)
T(s1, 0.5, 6.6, 4.0, 0.5, "Detect → Investigate manually → Guess → Decide", fs=12, bold=True, color=DARK)

# ==========================================
# SLIDE 2: SOLUTION
# ==========================================
R(s2, 0, 0, W, 0.8, DARK)
T(s2, 0.2, 0.1, 12, 0.4, "Trace turns KPI anomalies into evidence-backed decisions.", fs=20, bold=True, color=WHITE)
T(s2, 0.2, 0.45, 12, 0.3, "From signal → explanation → confidence → action", fs=14, color=RGBColor(0xCC, 0xCC, 0xFF))

stage_w = 2.4; space = 0.2; start_x = 0.2
stages = [
    ("1. DETECT SIGNAL", "Statistical baseline identifies meaningful deviation from normal behavior."),
    ("2. FIND THE DRIVER", "Decompose the KPI across dimensions to identify which segment actually explains the movement."),
    ("3. RETRIEVE EVIDENCE", "Pull relevant evidence from structured data + operational sources (GitHub, Jira, Slack)."),
    ("4. VALIDATE HYPOTHESES", "Connect temporal signals & supporting evidence, rank likely explanations, assign confidence."),
    ("5. RECOMMEND / ESCALATE", "High confidence → action\nLow confidence → investigate\nAmbiguous → competing explanations")
]
for i, (title, body) in enumerate(stages):
    x = start_x + i * (stage_w + space)
    R(s2, x, 1.1, stage_w, 2.0, LGRAY, border=MGRAY)
    R(s2, x, 1.1, stage_w, 0.4, PUR)
    T(s2, x+0.1, 1.2, stage_w-0.2, 0.3, title, fs=10, bold=True, color=WHITE)
    T(s2, x+0.1, 1.6, stage_w-0.2, 1.4, body, fs=9, color=DGRAY)
    if i < 4:
        RightArrow(s2, x+stage_w+0.05, 2.0, 0.1, 0.2, MGRAY)

# Left: TRACE IN ACTION
R(s2, 0.2, 3.4, 5.8, 0.4, DARK)
T(s2, 0.3, 3.5, 5.6, 0.3, "TRACE IN ACTION", fs=12, bold=True, color=WHITE)

# Before
R(s2, 0.2, 4.0, 1.6, 1.2, WHITE, border=MGRAY)
T(s2, 0.2, 4.2, 1.6, 0.3, "BEFORE", fs=11, bold=True, align=PP_ALIGN.CENTER, color=DGRAY)
T(s2, 0.2, 4.6, 1.6, 0.4, "Revenue ↓ 32%", fs=14, bold=True, align=PP_ALIGN.CENTER, color=RED)

RightArrow(s2, 1.9, 4.5, 0.2, 0.2, MGRAY)

# After Trace flow
x_after = 2.2; y_after = 4.0
R(s2, x_after, y_after, 3.8, 2.5, RGBColor(0xED, 0xE7, 0xF6), border=PUR)
T(s2, x_after+0.1, y_after+0.1, 3.6, 0.2, "AFTER TRACE", fs=10, bold=True, color=PUR)

flow_items = [
    "iOS → primary driver",
    "Checkout failures ↑",
    "Recent Stripe/payment deployment",
    "Supporting operational evidence"
]
curr_y = y_after + 0.35
for item in flow_items:
    T(s2, x_after+0.1, curr_y, 3.6, 0.2, item, fs=9.5, bold=True)
    if item != flow_items[-1]:
        Arrow(s2, x_after+1.8, curr_y+0.25, 0.1, 0.15, MGRAY)
    curr_y += 0.4

Arrow(s2, x_after+1.8, curr_y-0.1, 0.1, 0.15, PUR)

R(s2, x_after+0.1, curr_y+0.1, 3.6, 0.6, WHITE, border=PUR)
T(s2, x_after+0.15, curr_y+0.15, 3.3, 0.2, "LIKELY ROOT CAUSE (92% CONFIDENCE)", fs=9.5, bold=True, color=PUR)
T(s2, x_after+0.15, curr_y+0.4, 3.3, 0.2, "Recommended action: rollback / investigate deployment", fs=8.5, color=DARK)

# Right: MOST IMPORTANT DIFFERENTIATOR
R(s2, 6.4, 3.4, 6.7, 0.4, PUR)
T(s2, 6.5, 3.5, 6.5, 0.3, "MOST IMPORTANT DIFFERENTIATOR", fs=12, bold=True, color=WHITE)

R(s2, 6.4, 4.0, 6.7, 2.5, WHITE, border=MGRAY)
T(s2, 6.4, 4.3, 6.7, 0.4, "TRACE KNOWS WHEN IT DOESN’T KNOW.", fs=18, bold=True, align=PP_ALIGN.CENTER, color=DARK)

cw = 2.1
T(s2, 6.6, 5.0, cw, 0.3, "HIGH CONFIDENCE", fs=11, bold=True, color=GREEN)
T(s2, 6.6, 5.3, cw, 0.3, "→ Recommend action", fs=10, color=DGRAY)

T(s2, 8.8, 5.0, cw, 0.3, "LOW CONFIDENCE", fs=11, bold=True, color=RED)
T(s2, 8.8, 5.3, cw, 0.3, "→ Recommend investigation", fs=10, color=DGRAY)

T(s2, 10.9, 5.0, cw, 0.3, "AMBIGUOUS", fs=11, bold=True, color=AMBER)
T(s2, 10.9, 5.3, cw, 0.5, "→ Surface competing explanations\n→ No unsafe automated action", fs=10, color=DGRAY)

# FINAL TAKEAWAY ON SLIDE 2
R(s2, 0, 6.8, W, 0.7, DARK)
T(s2, 0.5, 7.0, W-1.0, 0.4, "Dashboards answer what changed. Trace answers why it changed, how sure we are, and what to do next.", fs=16, bold=True, align=PP_ALIGN.CENTER, color=WHITE)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
