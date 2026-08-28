"""
Trace.ai — AIC 2026 — 2-Slide Final PPT
Only 2 slides: Problem Statement Understanding & Proposed Solution
No mock data, explicitly mentions prototype is for video only.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = r'C:\Users\mohit\OneDrive\Desktop\AIC\ps\AIC_Talent-Brand_PPT-Template (1).pptx'
OUTPUT   = r'C:\Users\mohit\OneDrive\Desktop\AIC\TraceAI_Final_2Slides.pptx'

prs = Presentation(TEMPLATE)
W = 13.33; H = 7.50

PUR   = RGBColor(0x7B,0x2F,0xBE); DARK  = RGBColor(0x1A,0x1A,0x2E)
WHITE = RGBColor(0xFF,0xFF,0xFF); LGRAY = RGBColor(0xF5,0xF5,0xF5)
MGRAY = RGBColor(0xC8,0xC8,0xC8); DGRAY = RGBColor(0x50,0x50,0x50)
RED   = RGBColor(0xDC,0x26,0x26); GREEN = RGBColor(0x16,0xA3,0x4A)
AMBER = RGBColor(0xD9,0x7B,0x06); TEAL  = RGBColor(0x0D,0x9E,0x8E)

def I(v): return Inches(v)

def R(slide,l,t,w,h,fill,border=None):
    s=slide.shapes.add_shape(1,I(l),I(t),I(w),I(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if border: s.line.color.rgb=border; s.line.width=Pt(0.75)
    else: s.line.fill.background()
    return s

def T(slide,l,t,w,h,text,fs=9,bold=False,color=None,align=PP_ALIGN.LEFT,italic=False):
    color=color or DARK
    bx=slide.shapes.add_textbox(I(l),I(t),I(w),I(h))
    bx.text_frame.word_wrap=True; tf=bx.text_frame; tf.clear()
    p=tf.paragraphs[0]; p.alignment=align; run=p.add_run(); run.text=text
    run.font.size=Pt(fs); run.font.bold=bold; run.font.italic=italic
    run.font.color.rgb=color; return bx

def BL(slide,l,t,w,h,lines,fs=8.5,color=None,bullet="  - ",sp=2):
    color=color or DGRAY
    bx=slide.shapes.add_textbox(I(l),I(t),I(w),I(h))
    bx.text_frame.word_wrap=True; tf=bx.text_frame; first=True
    for line in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_before=Pt(sp); run=p.add_run(); run.text=bullet+line
        run.font.size=Pt(fs); run.font.color.rgb=color
    return bx

TABS = ["1. Problem Understanding","2. Proposed Solution"]
TW = 2.5

def footer(slide,page_num,active):
    R(slide,0,H-0.38,W,0.38,LGRAY,border=MGRAY)
    T(slide,0.1,H-0.34,0.28,0.28,str(page_num),fs=9,bold=True,color=DARK)
    for i,name in enumerate(TABS):
        x=0.42+i*TW
        bg=DARK if i==active else LGRAY; fg=WHITE if i==active else DGRAY
        R(slide,x,H-0.36,TW-0.06,0.3,bg)
        T(slide,x+0.04,H-0.34,TW-0.12,0.26,name,fs=7.5,bold=(i==active),color=fg,align=PP_ALIGN.CENTER)

def hdr(slide,title,sub):
    R(slide,0,0,W,0.6,DARK)
    T(slide,0.18,0.06,12.5,0.32,title,fs=16,bold=True,color=WHITE)
    T(slide,0.18,0.38,12.5,0.2,sub,fs=9,color=RGBColor(0xCC,0xCC,0xFF))

def sh(slide,l,t,w,h,text,fs=8.5):
    R(slide,l,t,w,h,PUR)
    T(slide,l+0.08,t+(h-0.2)/2,w-0.16,0.2,text,fs=fs,bold=True,color=WHITE)

def clear(slide):
    for s in slide.shapes:
        if s.has_text_frame: s.text_frame.clear()

# ══════════════════════════════════════════════
# SLIDE 1 — PROBLEM STATEMENT
# ══════════════════════════════════════════════
s4=prs.slides[3]; clear(s4)
hdr(s4,"Problem Statement Understanding: The Last-Mile Gap in Business Intelligence",
    "Modern BI tools are purely descriptive. They flag issues but fail to explain why they happened or what to do next.")

R(s4,0.18,0.72,W-0.36,0.38,RGBColor(0xED,0xE7,0xF6),border=PUR)
T(s4,0.28,0.76,W-0.56,0.3,
  '"A dashboard can show revenue dropped in a region; it rarely explains why or what to do next '
  '— that translation still falls to an analyst, often taking days."  — AIC Problem Statement',
  fs=8.5,italic=True,color=PUR)

pains=[
    ("TIME & REVENUE LOST",
     ["By the time an analyst manually diagnoses a KPI drop, gathers evidence, and writes a report, critical business hours (and revenue) have been lost.", "Manual diagnosis is inherently slow and reactive."]),
    ("CONTEXT BLINDNESS",
     ["Data is siloed. A drop in a business metric (like checkout conversions) cannot be easily correlated by a BI tool to an external technical event (like a bad code deployment).", "Lack of cross-functional visibility."]),
    ("LACK OF PRESCRIPTION",
     ["Existing tools alert users that something went wrong, but fail to provide ranked hypotheses or actionable remediation steps.", "Forces users to guess the root cause."]),
]
CW=4.22; CH=2.0; sy=1.22
for i,(label,bullets) in enumerate(pains):
    x=0.18+i*(CW+0.07)
    R(s4,x,sy,CW,CH,WHITE,border=MGRAY)
    R(s4,x,sy,CW,0.3,DARK)
    T(s4,x+0.1,sy+0.04,CW-0.2,0.22,label,fs=8,bold=True,color=WHITE)
    BL(s4,x+0.1,sy+0.4,CW-0.2,1.5,bullets,fs=9.5,bullet="  . ",sp=8)

sh(s4,0.18,3.4,W-0.36,0.26,"Current Manual Pipeline: Slow and Siloed",fs=8.5)
steps=[("KPI Drops",None),("Analyst Alerted",None),("Manual Log Search\nAcross Silos",None),
       ("Spreadsheets &\nDashboards Analyzed",None),("Hypothesis\nFormulated",None),("Report Filed\n(Days Later)",RED)]
SW=(W-0.36)/len(steps)
for i,(label,col) in enumerate(steps):
    x=0.18+i*SW; bg=col or LGRAY; fg=WHITE if col else DARK
    R(s4,x+0.04,3.72,SW-0.1,0.72,bg,border=MGRAY)
    T(s4,x+0.06,3.78,SW-0.16,0.62,label,fs=8.5,bold=True,color=fg,align=PP_ALIGN.CENTER)
    if i<len(steps)-1:
        T(s4,x+SW-0.12,3.96,0.2,0.3,">",fs=12,bold=True,color=PUR,align=PP_ALIGN.CENTER)

R(s4,0.18,4.7,W-0.36,1.2,WHITE,border=MGRAY)
R(s4,0.18,4.7,W-0.36,0.3,PUR)
T(s4,0.28,4.75,W-0.56,0.2, "Our Understanding of the Challenge", fs=9, bold=True, color=WHITE)
BL(s4,0.28,5.1,W-0.56,0.7, 
   ["The core challenge is to automate the manual translation layer between raw data anomalies and actionable business intelligence.",
    "We need an engine that bridges the gap in real-time, removing the reliance on manual cross-referencing."], 
   fs=10, sp=6)

footer(s4,1,0)

# ══════════════════════════════════════════════
# SLIDE 2 — PROPOSED SOLUTION
# ══════════════════════════════════════════════
s5=prs.slides[4]; clear(s5)
hdr(s5,"Proposed Solution: Trace.ai Architecture & Pipeline",
    "A 4-Stage KPI Storytelling Engine that automates root cause analysis and impact quantification.")

sh(s5,0.18,0.72,W-0.36,0.26,"4-Stage Solution Architecture",fs=8.5)
pipeline=[
    ("STEP 1","Bayesian Anomaly\nDetection","Continuously monitors KPI streams\nusing statistical models to filter out\nseasonal noise and flag true deviations."),
    ("STEP 2","Deterministic\nDecomposition","Algorithmically drills down into the\nanomaly across dimensions to isolate\nthe exact structural sub-segment."),
    ("STEP 3","Time-Bounded\nRAG Evidence","Automatically queries internal systems\nbounded by the exact time window of\nthe anomaly to gather factual evidence."),
    ("STEP 4","Structured\nHypothesis Engine","Passes context to an LLM which outputs\nranked hypotheses, confidence scores,\nand specific prescribed actions."),
]
PW3=(W-0.36)/4
for i,(num,title,body) in enumerate(pipeline):
    x=0.18+i*PW3; cw=PW3-0.08
    R(s5,x+0.04,1.04,cw,1.55,WHITE,border=MGRAY)
    R(s5,x+0.04,1.04,cw,0.3,PUR)
    T(s5,x+0.1,1.07,cw-0.12,0.22,num,fs=8,bold=True,color=WHITE)
    T(s5,x+0.1,1.38,cw-0.12,0.42,title,fs=10,bold=True,color=PUR)
    T(s5,x+0.1,1.84,cw-0.12,0.68,body,fs=8.5,color=DGRAY)
    if i<3:
        T(s5,x+PW3-0.1,1.6,0.2,0.3,">",fs=13,bold=True,color=PUR,align=PP_ALIGN.CENTER)

sh(s5,0.18,2.74,6.55,0.26,"Key Solution Capabilities",fs=8.5)
diffs=[
    ("Root Cause Contribution Trees",
     "Visually maps the chain of causality from the top-level metric drop down to the specific failing segment to provide immediate context."),
    ("Automated Impact Quantification",
     "Estimates the real-time financial cost of the anomaly to help executives prioritize incidents based on business impact."),
    ("Alternative Hypotheses Generation",
     "Mimics a senior analyst by providing multiple ranked root causes and remediation steps, rather than a single rigid answer."),
]
for i,(title,body) in enumerate(diffs):
    y=3.06+i*1.1
    R(s5,0.18,y,6.55,1.0,WHITE,border=MGRAY)
    R(s5,0.18,y,0.12,1.0,PUR)
    T(s5,0.38,y+0.08,6.2,0.26,title,fs=10,bold=True,color=DARK)
    T(s5,0.38,y+0.38,6.2,0.54,body,fs=9,color=DGRAY)

sh(s5,6.84,2.74,6.3,0.26,"Implementation & Pitch Video Prototype",fs=8.5)
R(s5,6.84,3.06,6.3,3.3,WHITE,border=MGRAY)

BL(s5,7.0,3.2,6.0,2.8,
   ["The Concept in Action: To demonstrate the feasibility of our planned solution architecture, we have built a functional frontend prototype for our pitch video. This prototype visualizes how the final product will look and operate.",
    "Tech Stack Plan: React (Frontend), FastAPI (Backend Pipeline), LLM APIs (Hypothesis Generation), Vector DB (Log Retrieval).",
    "Outcome: Reduces the time-to-root-cause from days to seconds, directly addressing the core Accenture problem statement."],
   fs=10.5, sp=14, bullet="  > ")

R(s5,0.18,6.58,W-0.36,0.36,PUR)
T(s5,0.28,6.62,W-0.56,0.28,
  "Trace.ai closes the BI last-mile gap: from anomaly detection to prescribed action in seconds.",
  fs=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
footer(s5,2,1)

slides_to_keep = [s4, s5]
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    if sld not in [slide._element for slide in slides_to_keep]:
        xml_slides.remove(sld)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
