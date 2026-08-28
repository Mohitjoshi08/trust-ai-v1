import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = r'C:\Users\mohit\OneDrive\Desktop\AIC\ps\AIC_Talent-Brand_PPT-Template (1).pptx'
OUTPUT = r'C:\Users\mohit\OneDrive\Desktop\AIC\TraceAI_Strict_2_Pages.pptx'

prs = Presentation(TEMPLATE)

# Delete ALL existing slides to guarantee EXACTLY 2 pages
xml_slides = prs.slides._sldIdLst
for sld in list(xml_slides):
    xml_slides.remove(sld)

# Add exactly 2 blank slides
blank_layout = prs.slide_layouts[6]
s1 = prs.slides.add_slide(blank_layout)
s2 = prs.slides.add_slide(blank_layout)

W = 13.33; H = 7.50
PUR = RGBColor(0x7B,0x2F,0xBE)
DARK = RGBColor(0x1A,0x1A,0x2E)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LGRAY = RGBColor(0xF5,0xF5,0xF5)
MGRAY = RGBColor(0xC8,0xC8,0xC8)
DGRAY = RGBColor(0x50,0x50,0x50)

def I(v): return Inches(v)

def R(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, fs=10, bold=False, color=None, align=PP_ALIGN.LEFT):
    color = color or DARK
    bx = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    bx.text_frame.word_wrap = True
    tf = bx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = color
    return bx

def BL(slide, l, t, w, h, lines, fs=10, color=None, bullet="• "):
    color = color or DGRAY
    bx = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    bx.text_frame.word_wrap = True
    tf = bx.text_frame
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = bullet + line
        run.font.size = Pt(fs)
        run.font.color.rgb = color
    return bx

# ==========================================
# PAGE 1: PROBLEM STATEMENT
# ==========================================
R(s1, 0, 0, W, 0.8, DARK)
T(s1, 0.4, 0.15, 12, 0.5, "Problem Statement: The Last-Mile Gap in Business Intelligence", fs=24, bold=True, color=WHITE)

R(s1, 0.4, 1.2, W-0.8, 1.2, LGRAY, border=PUR)
T(s1, 0.6, 1.4, W-1.2, 0.8, 
  "Modern BI tools are purely descriptive. They flag that an issue occurred but fail to explain WHY it happened or WHAT to do next. The translation from raw anomaly to actionable insight still falls to an analyst, often taking days.",
  fs=14, bold=True, color=PUR)

T(s1, 0.4, 2.8, 12, 0.4, "Key Pain Points & Inefficiencies", fs=18, bold=True, color=DARK)

col_w = (W - 1.2) / 3
pains = [
    ("Time & Revenue Lost", "Manual diagnosis is inherently slow. By the time an analyst diagnoses a KPI drop, gathers evidence, and writes a report, critical business hours and revenue have been permanently lost."),
    ("Context Blindness", "Data is completely siloed. A drop in a business metric cannot be easily correlated by a BI tool to an external technical event (like a code deployment) because the evidence lives in different systems."),
    ("Lack of Prescription", "Existing dashboard tools act only as alarms. They alert users that something went wrong, but force users to manually guess the root cause without providing actionable remediation steps.")
]

for i, (title, body) in enumerate(pains):
    x = 0.4 + i * (col_w + 0.2)
    R(s1, x, 3.4, col_w, 3.2, WHITE, border=MGRAY)
    R(s1, x, 3.4, col_w, 0.6, PUR)
    T(s1, x+0.1, 3.55, col_w-0.2, 0.4, title, fs=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s1, x+0.2, 4.3, col_w-0.4, 2.0, body, fs=12, color=DGRAY)


# ==========================================
# PAGE 2: PROPOSED SOLUTION
# ==========================================
R(s2, 0, 0, W, 0.8, DARK)
T(s2, 0.4, 0.15, 12, 0.5, "Proposed Solution: Trace.ai Architecture", fs=24, bold=True, color=WHITE)

T(s2, 0.4, 1.2, 6.0, 0.4, "4-Stage Engine Architecture", fs=18, bold=True, color=DARK)

stages = [
    ("1. Bayesian Anomaly Detection", "Continuously monitors KPI streams using statistical models to filter out seasonal noise and flag mathematically significant deviations."),
    ("2. Deterministic Decomposition", "Algorithmically drills down into the anomaly across dimensions (Region, Platform, Device) to isolate the exact structural sub-segment."),
    ("3. Time-Bounded RAG Evidence", "Automatically queries internal systems bounded by the exact time window of the anomaly to gather factual evidence."),
    ("4. Structured Hypothesis Engine", "Passes the context to an LLM which outputs ranked hypotheses, confidence scores, and specific prescribed actions.")
]

for i, (title, body) in enumerate(stages):
    y = 1.8 + i * 1.1
    R(s2, 0.4, y, 6.0, 0.9, WHITE, border=MGRAY)
    T(s2, 0.5, y+0.1, 5.8, 0.3, title, fs=12, bold=True, color=PUR)
    T(s2, 0.5, y+0.4, 5.8, 0.4, body, fs=10, color=DGRAY)

T(s2, 6.8, 1.2, 6.0, 0.4, "Solution Capabilities", fs=18, bold=True, color=DARK)

caps = [
    ("Automated Root Cause Mapping", "Visually maps the chain of causality from the top-level metric drop down to the specific failing segment."),
    ("Impact Quantification", "Estimates the financial cost of the anomaly to help prioritize incidents based on business impact."),
    ("Alternative Hypotheses", "Provides multiple ranked root causes and remediation steps rather than a single rigid answer.")
]

for i, (title, body) in enumerate(caps):
    y = 1.8 + i * 0.9
    R(s2, 6.8, y, 6.1, 0.7, LGRAY, border=MGRAY)
    T(s2, 6.9, y+0.1, 5.9, 0.25, title, fs=11, bold=True, color=DARK)
    T(s2, 6.9, y+0.35, 5.9, 0.3, body, fs=10, color=DGRAY)

R(s2, 6.8, 4.8, 6.1, 1.9, WHITE, border=PUR)
T(s2, 6.9, 5.0, 5.9, 0.3, "Implementation & Pitch Video Prototype", fs=14, bold=True, color=PUR)
T(s2, 6.9, 5.5, 5.9, 1.0, 
  "To demonstrate the feasibility of our planned solution architecture, we built a functional frontend prototype for our pitch video. This prototype visualizes how the final product will operate.\n\n"
  "It proves that we can reduce time-to-root-cause from days to seconds, directly addressing the core problem.",
  fs=11, color=DARK)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
